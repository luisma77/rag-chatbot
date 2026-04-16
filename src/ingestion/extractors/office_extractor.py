from pptx import Presentation
from openpyxl import load_workbook

from src.logger import get_logger

logger = get_logger(__name__)


def extract_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
        if parts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(parts))
    text = "\n\n".join(slides_text)
    logger.debug(f"PPTX '{file_path}': {len(prs.slides)} slides, {len(text)} chars")
    return text


def extract_xlsx(file_path: str) -> str:
    wb = load_workbook(file_path, read_only=True, data_only=True)
    all_text = []
    for sheet in wb.worksheets:
        rows_text = []
        for row in sheet:
            cells = [str(cell.value) for cell in row if cell.value is not None]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            all_text.append(f"[Hoja: {sheet.title}]\n" + "\n".join(rows_text))
    text = "\n\n".join(all_text)
    logger.debug(f"XLSX '{file_path}': {len(wb.worksheets)} hojas, {len(text)} chars")
    return text
